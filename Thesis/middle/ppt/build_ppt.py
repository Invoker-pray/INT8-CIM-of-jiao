"""
生成中期答辩 PPT
仿照 Thesis/start/start.pptx 风格
"""

import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── 路径 ──────────────────────────────────────────────────────────────
ROOT   = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
SRC    = os.path.join(ROOT, "Thesis/start/start.pptx")
FIG    = os.path.join(ROOT, "Thesis/middle/paper/figures")
FIG2   = os.path.join(ROOT, "Thesis/middle/paper/fig")
OUT    = os.path.join(ROOT, "Thesis/middle/ppt/middle.pptx")

os.makedirs(os.path.dirname(OUT), exist_ok=True)

# ── 颜色常量 ─────────────────────────────────────────────────────────
PURPLE = RGBColor(0x6D, 0x28, 0xD9)
NAVY   = RGBColor(0x2D, 0x38, 0x80)
DARK   = RGBColor(0x11, 0x18, 0x27)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xF3, 0xF4, 0xF6)

# ── 加载原版 PPT（取装饰图形） ───────────────────────────────────────
src_prs = Presentation(SRC)
src_slides = list(src_prs.slides)   # 0-based index

# ── 新建演示文稿，复用尺寸 ──────────────────────────────────────────
prs = Presentation()
prs.slide_width  = src_prs.slide_width    # 18288000 emu ≈ 20"
prs.slide_height = src_prs.slide_height   # 10287000 emu ≈ 11"

W = prs.slide_width
H = prs.slide_height
inch = 914400   # 1 inch in emu


def inches(n): return int(n * inch)


# ── 工具：从源 slide 复制指定 shape (by name) ──────────────────────
def copy_shapes_by_name(src_slide, dst_slide, *names):
    sp_tree = dst_slide.shapes._spTree
    for shape in src_slide.shapes:
        if shape.name in names:
            sp_tree.append(copy.deepcopy(shape._element))


def copy_all_freeforms(src_slide, dst_slide):
    """复制 src_slide 中所有 Freeform (type 5) 到 dst_slide"""
    sp_tree = dst_slide.shapes._spTree
    for shape in src_slide.shapes:
        if shape.shape_type == 5:   # FREEFORM
            sp_tree.append(copy.deepcopy(shape._element))


# ── 工具：添加文字框 ─────────────────────────────────────────────────
def add_text(slide, text, left, top, width, height,
             font_name="Calibri", font_size=24, bold=False,
             color=DARK, align=PP_ALIGN.LEFT,
             word_wrap=True, italic=False):
    txb = slide.shapes.add_textbox(inches(left), inches(top),
                                   inches(width), inches(height))
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name  = font_name
    f.size  = Pt(font_size)
    f.bold  = bold
    f.color.rgb = color
    if italic:
        f.italic = True
    return txb


def add_text_block(slide, lines, left, top, width, height,
                   font_name="Calibri", font_size=22, color=DARK,
                   line_spacing=None):
    """添加多段文字框（每个元素为 (text, bold, size) 或 str）"""
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn as _qn
    txb = slide.shapes.add_textbox(inches(left), inches(top),
                                   inches(width), inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, sz = item, False, font_size
        else:
            text, bold, sz = item[0], item[1], item[2] if len(item)>2 else font_size
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        f = run.font
        f.name  = font_name
        f.size  = _Pt(sz)
        f.bold  = bold
        f.color.rgb = color
    return txb


# ── 工具：添加图片 ──────────────────────────────────────────────────
def add_pic(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        print(f"  [WARN] 图片不存在: {path}")
        return None
    if width and height:
        return slide.shapes.add_picture(path,
                                        inches(left), inches(top),
                                        inches(width), inches(height))
    elif width:
        return slide.shapes.add_picture(path,
                                        inches(left), inches(top),
                                        inches(width))
    elif height:
        return slide.shapes.add_picture(path,
                                        inches(left), inches(top),
                                        height=inches(height))
    else:
        return slide.shapes.add_picture(path, inches(left), inches(top))


# ── 工具：添加紫色标题栏（仿 start.pptx 的 section 页） ────────────
def add_section_bar(slide, title_text):
    """左侧竖紫色色块 + 白色标题（用于 section 过渡页）"""
    from pptx.util import Pt
    bar = slide.shapes.add_shape(1,   # MSO_SHAPE_TYPE.RECTANGLE
                                  0, 0, inches(6), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    # 标题
    txb = slide.shapes.add_textbox(inches(0.4), inches(3),
                                    inches(5.2), inches(4))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    f = run.font
    f.name  = "Cormorant Garamond Bold Italics"
    f.size  = Pt(48)
    f.bold  = True
    f.color.rgb = WHITE
    return txb


# ── 工具：内容页标题 ─────────────────────────────────────────────────
def slide_title(slide, text, font_size=34):
    return add_text(slide, text, 0.3, 0.0, 17.0, 0.9,
                    font_name="Calibri", font_size=font_size,
                    bold=True, color=DARK)


# ── 空白 layout ──────────────────────────────────────────────────────
blank_layout = prs.slide_layouts[6]   # Blank


def new_slide():
    return prs.slides.add_slide(blank_layout)


# ════════════════════════════════════════════════════════════════════
# Slide 1 — 封面
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_all_freeforms(src_slides[0], s)

add_text(s, "面向边缘智能应用的低功耗 Int8\n数字存算一体架构设计及FPGA验证",
         2.0, 2.0, 12.0, 3.5,
         font_name="Calibri", font_size=52, bold=True, color=PURPLE,
         word_wrap=True)
add_text(s, "中期答辩", 2.0, 5.5, 8.0, 0.8,
         font_name="Calibri", font_size=36, bold=True, color=NAVY)
add_text(s, "微电子强芯 2201   焦宏宝   2221411116",
         2.0, 7.0, 11.0, 0.7,
         font_name="Segoe Print", font_size=28, color=NAVY)
add_text(s, "指导老师：梁峰",
         2.0, 8.0, 8.0, 0.7,
         font_name="Segoe Print", font_size=28, color=NAVY)
add_text(s, "2026 / 04 / 12",
         2.0, 9.0, 8.0, 0.7,
         font_name="Segoe Print", font_size=28, color=NAVY)

# ════════════════════════════════════════════════════════════════════
# Slide 2 — 目录
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_all_freeforms(src_slides[1], s)

add_text(s, "Agenda Overview", 3.0, 0.1, 14.0, 1.0,
         font_name="Cormorant Garamond Bold Italics",
         font_size=66, bold=True, color=NAVY)

items_left = [
    ("01", "研究背景与选题动机"),
    ("02", "总体系统架构"),
    ("03", "CIM 阵列硬件设计"),
    ("04", "INT8 量化与流水线"),
]
items_right = [
    ("05", "权重 SRAM 工程设计"),
    ("06", "PicoRV32 软核集成"),
    ("07", "FPGA 验证与性能评估"),
    ("08", "工作进展与后续计划"),
]

for i, (num, label) in enumerate(items_left):
    y = 2.8 + i * 1.5
    add_text(s, num, 3.0, y, 0.9, 0.7,
             font_name="Cormorant Garamond Bold Italics",
             font_size=36, bold=True, color=PURPLE)
    add_text(s, label, 4.1, y, 5.5, 0.7,
             font_name="Calibri", font_size=28, color=DARK)

for i, (num, label) in enumerate(items_right):
    y = 2.8 + i * 1.5
    add_text(s, num, 10.5, y, 0.9, 0.7,
             font_name="Cormorant Garamond Bold Italics",
             font_size=36, bold=True, color=PURPLE)
    add_text(s, label, 11.6, y, 6.5, 0.7,
             font_name="Calibri", font_size=28, color=DARK)


# ════════════════════════════════════════════════════════════════════
# Slide 3 — 研究背景：内存墙问题
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[2], s, "Freeform 4")

slide_title(s, "01  研究背景：内存墙问题与边缘 AI 瓶颈")

add_pic(s, os.path.join(FIG, "m1.png"), 0.2, 1.0, 9.0, 9.5)

add_text_block(s, [
    ("边缘 AI 推理场景", True, 26),
    ("YOLO / SSD 等目标检测算法计算密集，实时性要求高", False, 22),
    ("电池供电、散热受限，功耗约束严苛", False, 22),
    ("", False, 10),
    ("冯·诺依曼架构瓶颈", True, 26),
    ("处理器与存储器分离，权重/激活频繁搬运", False, 22),
    ("访存能耗 ≈ 计算能耗 200×", False, 22),
    ("存储带宽增速远低于计算能力增速（内存墙）", False, 22),
    ("", False, 10),
    ("CIM 的核心思路", True, 26),
    ('在存储阵列内部完成计算，将"搬运"变为"原位计算"', False, 22),
    ("从体系结构层面缓解访存瓶颈", False, 22),
], 9.6, 1.1, 9.8, 9.5,
   font_name="Segoe Print", font_size=22, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 4 — 研究背景：CIM 技术全谱系
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[3], s, "Freeform 2")

slide_title(s, "01  研究背景：CIM 技术全谱系与本设计定位")

add_pic(s, os.path.join(FIG, "m2.png"), 0.2, 1.0, 10.0, 9.5)

add_text_block(s, [
    ("技术路线对比", True, 26),
    ("模拟 CIM（RRAM/PCM）：高并行，需 ADC/DAC，精度/噪声挑战大", False, 21),
    ("数字 SRAM-CIM：速度快、工艺成熟、数字设计流程兼容", False, 21),
    ("", False, 8),
    ("本课题定位", True, 26),
    ("数字 SRAM-CIM + Signed-INT8 + FPGA 验证", False, 21),
    ("目标：工程闭环，可复现的 RTL→仿真→上板验证", False, 21),
    ("", False, 8),
    ("趋势", True, 26),
    ("从模拟/新型器件主导 → 全数字并行发展", False, 21),
    ("从低比特 → 多比特 MAC（INT8/BF16）", False, 21),
    ("从 CNN → Transformer / 多模态", False, 21),
    ("需要体系结构层面可复现评估", False, 21),
], 10.5, 1.1, 9.0, 9.5,
   font_name="Segoe Print", font_size=21, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 5 — 两种实现极端对比（节标题 + 表格文字）
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[4], s, "Freeform 13", "Freeform 14")

slide_title(s, "01  两种数字 SRAM-CIM 实现极端与本课题取舍")

# 表头
headers = ["维度", "极端一：bit-serial popcount", "极端二：字级行为 MAC（本课题）"]
col_x   = [0.3, 4.8, 11.5]
col_w   = [4.3, 6.5, 8.0]

rows = [
    ("存储单元粒度",       "1-bit cell，AND + popcount",     "INT8/INT32 字，BRAM"),
    ("INT8 MAC 周期数",    "8 拍（bit-plane shift-add）",    "1 拍（组合乘加）"),
    ("阵列规模示例",       "128×128 位 × 8 bit-plane",       "16×16 INT8 MAC"),
    ("FPGA 资源映射",      "LUT+FF 为主，BRAM/DSP 低",       "BRAM bank + DSP48 推断"),
    ("精度后处理位置",     "CPU（软件完成）",                 "硬件流水内完成"),
    ('"存算一体"严格性',   "强（cell 即 MAC）",               "弱（近存储字级 MAC）"),
    ("FPGA 原型工程性",    "中",                              "高 ✓"),
    ("本课题选择",         "—",                              "✓ PYNQ-Z2 资源约束下最优取舍"),
]

y0 = 1.1
row_h = 0.88

# 画表头
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    cell = s.shapes.add_shape(1, inches(cx), inches(y0),
                              inches(cw - 0.05), inches(row_h))
    cell.fill.solid()
    cell.fill.fore_color.rgb = PURPLE
    cell.line.fill.background()
    txb = s.shapes.add_textbox(inches(cx + 0.1), inches(y0 + 0.12),
                                inches(cw - 0.2), inches(row_h - 0.15))
    tf = txb.text_frame
    p  = tf.paragraphs[0]
    run = p.add_run()
    run.text = hdr
    run.font.name  = "Calibri"
    run.font.size  = Pt(20)
    run.font.bold  = True
    run.font.color.rgb = WHITE

# 画数据行
for ri, row in enumerate(rows):
    y = y0 + (ri + 1) * row_h
    bg = GRAY if ri % 2 == 0 else WHITE
    for j, (cell_text, cx, cw) in enumerate(zip(row, col_x, col_w)):
        rect = s.shapes.add_shape(1, inches(cx), inches(y),
                                  inches(cw - 0.05), inches(row_h - 0.03))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb      = RGBColor(0xD1, 0xD5, 0xDB)
        rect.line.width          = Pt(0.5)
        fc = NAVY if j > 0 else DARK
        txb = s.shapes.add_textbox(inches(cx + 0.1), inches(y + 0.1),
                                    inches(cw - 0.2), inches(row_h - 0.15))
        tf = txb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        run = p.add_run()
        run.text = cell_text
        run.font.name  = "Calibri"
        run.font.size  = Pt(18)
        run.font.bold  = (ri == len(rows) - 1)
        run.font.color.rgb = fc


# ════════════════════════════════════════════════════════════════════
# Slide 6 — 总体系统架构
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[7], s, "Freeform 32", "Freeform 33")

slide_title(s, "02  总体系统架构")

add_pic(s, os.path.join(FIG, "m5.png"), 0.2, 1.0, 12.5, 9.5)

add_text_block(s, [
    ("两种控制模式", True, 26),
    ("ARM PS 模式：Python/MMIO 驱动", False, 22),
    ("PicoRV32 模式：纯 PL 自主推理", False, 22),
    ("", False, 8),
    ("核心数据通路", True, 26),
    ("权重 DMA 加载 → 推理触发", False, 22),
    ("INT8 输入 → MVM → Bias → ReLU", False, 22),
    ("重量化 → INT8 输出", False, 22),
    ("", False, 8),
    ("目标平台", True, 26),
    ("PYNQ-Z2 (Zynq-7020)", False, 22),
    ("60 MHz，220 DSP48", False, 22),
], 12.9, 1.1, 6.8, 9.5,
   font_name="Segoe Print", font_size=22, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 7 — CIM Tile 16×16 设计
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[8], s, "Freeform 2", "Freeform 4")

slide_title(s, "03  CIM 阵列：16×16 CIM Tile 设计")

add_pic(s, os.path.join(FIG, "m3.png"), 0.2, 1.1, 11.5, 9.4)

add_text_block(s, [
    ("核心计算单元", True, 26),
    ("16×16 纯组合逻辑 MAC 阵列", False, 22),
    ("一拍完成 256 次 INT8 × INT8 = INT32", False, 22),
    ("", False, 8),
    ("输入广播机制", True, 26),
    ("x_eff[c]（9-bit UINT）按行广播", False, 22),
    ("x_eff = pixel - input_zp（默认−128）", False, 22),
    ("范围 [128, 383] 确保非负", False, 22),
    ("", False, 8),
    ("列方向链式累加", True, 26),
    ("row_acc[c+1] = row_acc[c] + x×w", False, 22),
    ("输出 16 个 INT32 部分和", False, 22),
    ("", False, 8),
    ("为何选纯组合设计", True, 26),
    ('时序控制集中在 FSM，Tile 做"存储做计算"的数字抽象', False, 22),
    ("DSP48E1 自动推断，综合器友好", False, 22),
], 12.0, 1.1, 7.6, 9.5,
   font_name="Segoe Print", font_size=22, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 8 — INT8 量化数据通路
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[9], s, "Freeform 5")

slide_title(s, "04  INT8 量化与重定标数据通路")

add_pic(s, os.path.join(FIG, "m5.png"), 0.2, 1.1, 11.0, 9.4)

add_text_block(s, [
    ("量化链路（5 步）", True, 26),
    ("① 零点减法  x_eff = pixel − zp", False, 22),
    ("② CIM Tile  psum = Σ x_eff × w", False, 22),
    ("③ 偏置累加  psum += bias[r]（INT32）", False, 22),
    ("④ ReLU 激活（ACT_MODE via CSR）", False, 22),
    ("⑤ 重量化：× requant_mult → >> shift → clamp", False, 22),
    ("", False, 8),
    ("位宽演变", True, 26),
    ("UINT8 → 9b → 17b(积) → 32b(累加)", False, 22),
    ("→ 64b(重量化乘) → 32b(右移) → INT8", False, 22),
    ("", False, 8),
    ("硬件端到端闭环", True, 26),
    ("PS 仅在层间介入，单层推理完全在 PL", False, 22),
    ("全程与 Python golden model bit-exact", False, 22),
], 11.4, 1.1, 8.2, 9.5,
   font_name="Segoe Print", font_size=22, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 9 — 7 级流水线 FSM
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[8], s, "Freeform 2", "Freeform 4")

slide_title(s, "04  7 级流水线状态机（17 状态 FSM）")

add_pic(s, os.path.join(FIG, "m6.png"), 0.2, 1.1, 11.0, 9.4)

add_text_block(s, [
    ("计算流水线（每 IB 迭代）", True, 25),
    ("ST_FETCH / ST_WAIT_SRAM  发出 BRAM 读地址", False, 21),
    ("ST_XEFF_REG  零点减法→锁存 x_eff  ~10 ns", False, 21),
    ("ST_MAC  16 元素 MAC 链→锁存 psum  ~10 ns", False, 21),
    ("ST_COMPUTE  psum_accum += tile_psum  ~4 ns", False, 21),
    ("", False, 6),
    ("输出流水线（每 OB 完成后）", True, 25),
    ("ST_BIAS_ADD → ST_ACTIVATE  加偏置 / ReLU", False, 21),
    ("ST_STORE  64-bit 重量化乘法", False, 21),
    ("ST_SHIFT_CLAMP → ST_WRITE_OBUF  写输出缓冲", False, 21),
    ("", False, 6),
    ("时序优化四阶段", True, 25),
    ("25 MHz → patch1 流水线拆分 → 40 MHz", False, 21),
    ("→ patch2 计算段三分 → 50 MHz", False, 21),
    ("→ patch3 ST_STORE 再拆 → 60 MHz ✓", False, 21),
    ("关键路径：w_tile_reg→DSP48→CARRY4→tile_psum_reg", False, 19),
], 11.5, 1.1, 8.1, 9.5,
   font_name="Segoe Print", font_size=21, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 10 — 权重 SRAM 16-bank 方案
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[3], s, "Freeform 2")

slide_title(s, "05  权重 SRAM：16-Bank 拆分与 BRAM 推断")

add_pic(s, os.path.join(FIG, "m4.png"), 0.2, 1.1, 11.5, 9.4)

add_text_block(s, [
    ("问题：BRAM 推断失败", True, 26),
    ("原设计 2048-bit 整块 SRAM + bit-select 写入", False, 22),
    ("Vivado 放弃 BRAM→退化为寄存器（资源爆炸）", False, 22),
    ("", False, 8),
    ("解决方案：16 独立 Bank", True, 26),
    ("每个 Bank：128-bit × 392-depth", False, 22),
    ("对应 Tile 一行权重（TILE_COLS×WEIGHT_W）", False, 22),
    ("generate 实例化，每个 Bank 独立推断 BRAM", False, 22),
    ("", False, 8),
    ("写入路径：Read-Modify-Write", True, 26),
    ("4 次 32-bit AXI 写入 → staging register 累积", False, 22),
    ("满 128-bit → 单周期 whole-word 写回 Bank", False, 22),
    ('确保 BRAM 推断所需"纯整字写入"', False, 22),
    ("", False, 8),
    ("读取：16 Bank 并行输出 → 2048-bit → CIM Tile", False, 22),
], 12.0, 1.1, 7.6, 9.5,
   font_name="Segoe Print", font_size=22, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 11 — AXI4-Lite 接口 + 时序收敛
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[4], s, "Freeform 13", "Freeform 14")

slide_title(s, "05  AXI4-Lite 接口与时序收敛迭代")

# 左侧 AXI 地址映射表
add_text(s, "CSR 地址映射（14-bit，16 KB）", 0.3, 1.0, 9.0, 0.7,
         font_name="Calibri", font_size=24, bold=True, color=DARK)

axi_rows = [
    ("0x000–0x00C", "控制/状态/IRQ 寄存器"),
    ("0x010–0x02C", "层配置（IN_DIM, OUT_DIM, 量化参数）"),
    ("0x030–0x03C", "性能计数器（cycle / MAC count）"),
    ("0x040",       "argmax 预测结果"),
    ("0x044–0x04C", "权重 DMA（WDMA_ADDR / DATA / CTRL）"),
    ("0x100–0x2FF", "Logits 回读窗口"),
    ("0x1000–0x1FFF", "输入缓冲写入窗口"),
    ("0x2000–0x2FFF", "偏置缓冲写入窗口"),
]
for ri, (addr, desc) in enumerate(axi_rows):
    y = 1.8 + ri * 0.88
    bg = GRAY if ri % 2 == 0 else WHITE
    for j, (txt, cx, cw) in enumerate([
            (addr, 0.3, 3.5), (desc, 3.9, 6.0)]):
        rect = s.shapes.add_shape(1, inches(cx), inches(y),
                                  inches(cw - 0.05), inches(0.83))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb      = RGBColor(0xD1, 0xD5, 0xDB)
        rect.line.width          = Pt(0.5)
        txb = s.shapes.add_textbox(inches(cx + 0.1), inches(y + 0.1),
                                    inches(cw - 0.2), inches(0.7))
        tf = txb.text_frame
        p  = tf.paragraphs[0]
        run = p.add_run()
        run.text = txt
        run.font.name  = "Consolas" if j == 0 else "Calibri"
        run.font.size  = Pt(18)
        run.font.color.rgb = NAVY if j == 0 else DARK

# 右侧时序收敛表
add_text(s, "时序收敛四阶段", 10.3, 1.0, 9.0, 0.7,
         font_name="Calibri", font_size=24, bold=True, color=DARK)
timing_rows = [
    ("初始版本",    "25 MHz",  "组合逻辑 26 ns，频率受限"),
    ("patch 1",     "40 MHz",  "bias/act/requant/store 四级流水"),
    ("patch 2",     "50 MHz",  "计算段三分：XEFF_REG + MAC + COMPUTE"),
    ("patch 3",     "60 MHz ✓","ST_STORE 再拆为三步；WNS +0.33 ns"),
    ("55 MHz 变体", "55 MHz",  "PAR_OB=1，资源更紧凑的替代版本"),
]
tcols = [("版本", 2.2), ("频率", 1.8), ("说明", 5.2)]
tx0 = [10.3, 12.6, 14.5]
for ri, row in enumerate(timing_rows):
    y = 1.8 + ri * 0.88
    bg = RGBColor(0xED, 0xE9, 0xFE) if "✓" in row[1] else (GRAY if ri % 2 == 0 else WHITE)
    for j, (txt, cx) in enumerate(zip(row, tx0)):
        cw = tcols[j][1]
        rect = s.shapes.add_shape(1, inches(cx), inches(y),
                                  inches(cw - 0.05), inches(0.83))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb      = RGBColor(0xD1, 0xD5, 0xDB)
        rect.line.width          = Pt(0.5)
        txb = s.shapes.add_textbox(inches(cx + 0.1), inches(y + 0.1),
                                    inches(cw - 0.2), inches(0.7))
        tf = txb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        run = p.add_run()
        run.text = txt
        run.font.name  = "Calibri"
        run.font.size  = Pt(18)
        run.font.bold  = ("✓" in row[1])
        run.font.color.rgb = PURPLE if "✓" in row[1] else DARK


# ════════════════════════════════════════════════════════════════════
# Slide 12 — PicoRV32 设计动机与系统架构
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[8], s, "Freeform 2", "Freeform 4")

slide_title(s, "06  PicoRV32 软核集成：设计动机与架构")

add_pic(s, os.path.join(FIG, "m8.png"), 0.2, 1.1, 11.5, 9.4)

add_text_block(s, [
    ("设计动机", True, 26),
    ("任务书要求：软核替代 ARM PS", False, 22),
    ("Step 1-3 是 ARM Python/MMIO 驱动", False, 22),
    ("目标：纯 PL 侧自主推理，不依赖 PS", False, 22),
    ("", False, 8),
    ("为何选 PicoRV32", True, 26),
    ("MIT License，单文件 picorv32.v", False, 22),
    ("RV32IM，~1500 LUT", False, 22),
    ("原生 Memory Interface，易桥接 CSR", False, 22),
    ("", False, 8),
    ("系统架构", True, 26),
    ("PicoRV32 → Bridge → FW BRAM", False, 22),
    ("→ CIM IP（完全不改 AXI slave）", False, 22),
    ("→ Result BRAM / UART TX", False, 22),
    ("PS 仅提供时钟 + AXI 固件加载", False, 22),
], 12.0, 1.1, 7.6, 9.5,
   font_name="Segoe Print", font_size=22, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 13 — PicoRV32 Bridge 与关键设计
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[4], s, "Freeform 13", "Freeform 14")

slide_title(s, "06  PicoRV32：Bus Bridge 与关键工程问题")

add_pic(s, os.path.join(FIG, "m9.png"), 0.2, 1.1, 9.5, 9.4)

add_text_block(s, [
    ("Wishbone → CSR Bridge", True, 26),
    ("addr[31:30] 两位解码四个外设", False, 22),
    ("00→FW BRAM, 01→CIM, 10→UART, 11→Result BRAM", False, 21),
    ("发出与 ARM 完全相同的 AXI 事务", False, 22),
    ("", False, 8),
    ("双端口 FW BRAM", True, 26),
    ("Port A：RV32 取指令 / 数据", False, 22),
    ("Port B：PS 通过 AXI 写入 firmware.hex", False, 22),
    ("运行时加载，无需重新生成 bitstream", False, 22),
    ("", False, 8),
    ("UART 问题 → Result BRAM 方案", True, 26),
    ("PL UART 与 PS USB 隔离，无法直接读取", False, 22),
    ("PicoRV32 写结果到 Result BRAM Port A", False, 22),
    ("PS 通过 AXI 读 Port B，控制 cpu_rst_n", False, 22),
    ("", False, 8),
    ("模型规模约束", True, 26),
    ("固件放入 32 KB BRAM：784→16→10，约 14.5 KB", False, 22),
    ("INT8 精度 90%+（量化损失，非硬件错误）", False, 22),
], 9.9, 1.1, 9.7, 9.5,
   font_name="Segoe Print", font_size=22, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 14 — 验证方法学
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[10], s, "Freeform 4")

slide_title(s, "07  验证方法学：三层一致性验证体系")

add_pic(s, os.path.join(FIG, "m10.png"), 0.2, 1.1, 11.5, 9.4)

add_text_block(s, [
    ("第一层：Python Golden Model", True, 24),
    ("bit-accurate INT8 整数推理", False, 21),
    ("生成 hex 测试向量和期望输出", False, 21),
    ("", False, 6),
    ("第二层：VCS RTL 仿真", True, 24),
    ("tb_cim_tile（103 随机用例）", False, 21),
    ("tb_cim_accel_core（边界/随机）", False, 21),
    ("tb_mnist_e2e（完整 MLP 端到端）", False, 21),
    ("run_regression.sh 一键回归", False, 21),
    ("", False, 6),
    ("第三层：PYNQ-Z2 板级验证", True, 24),
    ("MLP：20 张 100%，bit-exact 100%", False, 21),
    ("LeNet-5：200 张 99.5%，bit-exact 100%", False, 21),
    ("PicoRV32：200 张 96.5%，logits bit-exact", False, 21),
    ("", False, 6),
    ("关键原则", True, 24),
    ("三层逐元素比对，每层输出与 golden 一致", False, 21),
    ("Pytest 回归（16 测试，0.29 s，全 GREEN）", False, 21),
], 12.0, 1.1, 7.6, 9.5,
   font_name="Segoe Print", font_size=21, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 15 — FPGA 综合结果
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[4], s, "Freeform 13", "Freeform 14")

slide_title(s, "07  FPGA 综合结果：资源占用与时序")

add_pic(s, os.path.join(FIG, "m11.png"), 0.2, 1.1, 9.5, 9.4)

# 右侧两个表格
add_text(s, "ARM PS 模式（60 MHz）", 10.0, 1.1, 9.5, 0.7,
         font_name="Calibri", font_size=24, bold=True, color=DARK)

res_rows_arm = [
    ("LUT",    "11,087", "20.84%"),
    ("FF",     "5,385",  "5.06%"),
    ("BRAM",   "35 块",   "25.00%"),
    ("DSP48",  "220",    "100.00% ⚠"),
    ("动态功耗", "0.284 W", "CIM IP"),
    ("WNS",    "+0.33 ns", "时序收敛 ✓"),
]
for ri, (name, val, pct) in enumerate(res_rows_arm):
    y = 1.9 + ri * 0.82
    for j, (txt, cx, cw) in enumerate([
            (name, 10.0, 3.0), (val, 13.1, 2.8), (pct, 16.0, 3.5)]):
        bg = RGBColor(0xED, 0xE9, 0xFE) if "⚠" in pct or "✓" in pct else (
             GRAY if ri % 2 == 0 else WHITE)
        rect = s.shapes.add_shape(1, inches(cx), inches(y),
                                  inches(cw - 0.05), inches(0.77))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        rect.line.width = Pt(0.5)
        txb = s.shapes.add_textbox(inches(cx + 0.1), inches(y + 0.08),
                                    inches(cw - 0.2), inches(0.62))
        p = txb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = txt
        run.font.name  = "Calibri"
        run.font.size  = Pt(18)
        run.font.bold  = (j == 0)
        run.font.color.rgb = PURPLE if ("⚠" in pct or "✓" in pct) else DARK

add_text(s, "PicoRV32 模式（50 MHz）额外开销",
         10.0, 6.8, 9.5, 0.7,
         font_name="Calibri", font_size=22, bold=True, color=DARK)
add_text_block(s, [
    ("+1,627 LUT  /  +887 FF", False, 22),
    ("+9 BRAM（FW BRAM + Result BRAM）", False, 22),
    ("+0 DSP（RV32 不占用 DSP）", False, 22),
    ("logits 与 ARM 路径 bit-exact 100% ✓", True, 22),
], 10.0, 7.5, 9.5, 3.0,
   font_name="Segoe Print", font_size=22, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 16 — MLP 推理验证结果
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[4], s, "Freeform 13", "Freeform 14")

slide_title(s, "07  MLP 推理验证：784→128→10，MNIST")

add_text_block(s, [
    ("网络结构", True, 26),
    ("784 → 128 → 10（2 层全连接 + ReLU + 重量化）", False, 22),
    ("INT8 量化，模型精度约 93%（完整测试集）", False, 22),
    ("", False, 8),
    ("硬件性能（ARM 模式，60 MHz）", True, 26),
    ("FC1（784→128）：3,136 cycles，MAC 吞吐 31.2 MAC/cycle", False, 22),
    ("FC2（128→10）：146 cycles", False, 22),
    ("总计 3,282 cycles，推理延迟 54.7 µs", False, 22),
    ("", False, 8),
    ("验证结果", True, 26),
    ("20 张 MNIST 分类准确率：100% (20/20)", False, 22),
    ("与 Python golden 逐元素比对：bit-exact 100%", False, 22),
    ("", False, 8),
    ("加速比（vs ARM Cortex-A9 NumPy 基线）", True, 26),
    ("SW baseline：1,410.8 µs（NumPy @ 650 MHz）", False, 22),
    ("加速比：25.8×（HW @ 60 MHz vs SW @ 650 MHz）", True, 24),
    ("注：硬件时钟仅为软件的 1/10.8", False, 20),
], 0.3, 1.1, 9.5, 9.5,
   font_name="Segoe Print", font_size=22, color=NAVY)

add_pic(s, os.path.join(FIG, "m13.png"), 9.9, 1.1, 9.6, 9.4)


# ════════════════════════════════════════════════════════════════════
# Slide 17 — LeNet-5 CNN 验证结果
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[2], s, "Freeform 4")

slide_title(s, "07  LeNet-5 CNN 验证：7 层网络，200 张 MNIST")

add_pic(s, os.path.join(FIG, "m12.png"), 0.2, 1.1, 10.0, 9.4)

add_text_block(s, [
    ("网络结构（7 层）", True, 24),
    ("Conv1→Pool1→Conv2→Pool2→FC3→FC4→FC5", False, 20),
    ("im2col 展开 Conv → MVM 引擎，软硬件异构方案", False, 20),
    ("", False, 6),
    ("逐层周期数（60 MHz ARM 模式）", True, 24),
    ("Conv1（1×28×28→6×24×24）：63,360 cycles", False, 20),
    ("Conv2（6×12×12→16×8×8）：10,112 cycles", False, 20),
    ("FC3（256→120）：1,552 cycles", False, 20),
    ("FC4/FC5（120→84→10）：1,010 cycles", False, 20),
    ("总计：76,034 cycles，延迟 1,267.2 µs", False, 20),
    ("", False, 6),
    ("验证结果", True, 24),
    ("200 张准确率：99.5% (199/200)", True, 22),
    ("bit-exact 100%（含唯一误分类均为量化损失）", False, 20),
    ("", False, 6),
    ("加速比", True, 24),
    ("SW baseline（sequential）：26,466 µs", False, 20),
    ("加速比：20.9×", True, 22),
], 10.4, 1.1, 9.2, 9.5,
   font_name="Segoe Print", font_size=20, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 18 — SQ-mapping 优化 + 延迟分解
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[9], s, "Freeform 5")

slide_title(s, "07  SQ-mapping 软件优化与端到端延迟拆解")

add_pic(s, os.path.join(FIG2, "latency_breakdown.png"), 0.2, 1.1, 9.5, 9.4)

add_text_block(s, [
    ("SQ-mapping 权重打包优化", True, 26),
    ("当 col_len << MAX_IN_DIM 时，多像素并行拼包", False, 22),
    ("", False, 6),
    ("Conv1（col_len=25，C_out=6）", True, 22),
    ("打包系数 min(784/25, 128/6)=21", False, 20),
    ("784 次 MVM → 38 次  → 理论 ~20×", False, 20),
    ("", False, 6),
    ("Conv2（col_len=150，C_out=16）", True, 22),
    ("打包系数=5，100 次→20 次  → 理论 ~5×", False, 20),
    ("", False, 6),
    ("实测加速（LeNet-5，200 张）", True, 26),
    ("未打包：708.3 s（3.54 s/张）", False, 22),
    ("打包后：325.1 s（1.63 s/张）", False, 22),
    ("墙钟加速：2.18×", True, 24),
    ("", False, 6),
    ("瓶颈分析（延迟分解饼图）", True, 26),
    ("Conv1 HW 计算仅 4.1 ms", False, 22),
    ("MMIO 权重搬运 ~620 ms（主要瓶颈）", False, 22),
    ("→ 根本解法：AXI4-Full burst DMA（未来工作）", False, 22),
], 9.9, 1.1, 9.7, 9.5,
   font_name="Segoe Print", font_size=22, color=NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 19 — 两种控制模式对比
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[4], s, "Freeform 13", "Freeform 14")

slide_title(s, "07  两种控制模式综合对比")

add_pic(s, os.path.join(FIG, "m14.png"), 0.2, 1.1, 9.8, 9.4)

# 对比表
add_text(s, "指标对比", 10.3, 1.1, 9.0, 0.7,
         font_name="Calibri", font_size=24, bold=True, color=DARK)

cmp_rows = [
    ("指标",             "ARM PS 模式",    "PicoRV32 模式"),
    ("工作频率",          "60 MHz",         "50 MHz"),
    ("时序 WNS",         "−0.086 ns",      "+0.204 ns ✓"),
    ("LUT",              "11,087 (20.84%)", "12,714 (23.90%)"),
    ("BRAM",             "35 (25.00%)",    "44 (31.43%)"),
    ("动态功耗",          "1.807 W",        "1.745 W"),
    ("自主推理",          "需 PS Python",   "纯 PL ✓"),
    ("推理准确率",        "99.5%（LeNet）",  "96.5%（小 MLP）"),
    ("bit-exact",        "100%",           "100% ✓"),
]
for ri, row in enumerate(cmp_rows):
    y = 1.9 + ri * 0.88
    for j, (txt, cx, cw) in enumerate([
            (row[0], 10.3, 2.8),
            (row[1], 13.2, 3.0),
            (row[2], 16.3, 3.3)]):
        is_hdr = (ri == 0)
        bg = PURPLE if is_hdr else (
             RGBColor(0xED, 0xE9, 0xFE) if "✓" in row[2] else
             (GRAY if ri % 2 == 0 else WHITE))
        fc = WHITE if is_hdr else (PURPLE if "✓" in txt else DARK)
        rect = s.shapes.add_shape(1, inches(cx), inches(y),
                                  inches(cw - 0.05), inches(0.83))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        rect.line.width = Pt(0.5)
        txb = s.shapes.add_textbox(inches(cx + 0.08), inches(y + 0.1),
                                    inches(cw - 0.15), inches(0.7))
        p = txb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = txt
        run.font.name  = "Calibri"
        run.font.size  = Pt(18)
        run.font.bold  = is_hdr or "✓" in txt
        run.font.color.rgb = fc


# ════════════════════════════════════════════════════════════════════
# Slide 20 — 工作进展与后续计划
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_shapes_by_name(src_slides[7], s, "Freeform 32", "Freeform 33")

slide_title(s, "08  工作进展与后续计划")

# 左：已完成
add_text(s, "✅ 已完成工作", 0.3, 1.1, 9.0, 0.7,
         font_name="Calibri", font_size=26, bold=True, color=DARK)
done = [
    "Step 1  RTL 设计 + VCS 仿真（3 个 testbench + regression）",
    "Step 2  PYNQ-Z2 上板验证（MLP + LeNet-5 CNN）",
    "Step 3  多层自动推理驱动 + im2col Conv 支持",
    "Step 4  PicoRV32 集成，纯 PL 自主推理",
    "Step 6  SQ-mapping 打包优化（2.18× 加速）",
    "Step 6  逐层 bit-exact 验证基础设施",
    "Step 6  端到端延迟分解 profiler + 饼图",
    "Step 6  Pytest 回归（16 测试，全 GREEN）",
    "         中期论文撰写（共 36 页）",
]
for i, txt in enumerate(done):
    add_text(s, txt, 0.5, 1.9 + i * 0.9, 9.5, 0.85,
             font_name="Calibri", font_size=20, color=DARK)

# 右：后续计划
add_text(s, "📋 后续计划", 10.2, 1.1, 9.0, 0.7,
         font_name="Calibri", font_size=26, bold=True, color=DARK)

plan = [
    ("Step 5  KV260 移植（进行中）",
     "PAR_OB=4，UltraScale+ 200 MHz 目标"),
    ("Phase D  UltraRAM 替换 BRAM",
     "MAX_IN_DIM 从 784 扩展到 1024+"),
    ("Phase C3  AXI4-Full burst DMA",
     "消除 MMIO 搬运瓶颈（~700 ms/张），理论 100× 加速"),
    ("Phase B1  资源/时序自动提取",
     "build_history.csv，趋势线图表"),
    ("Phase E2  统一 CLI 入口",
     "sw/scripts/cim.py run，答辩 demo 干净"),
    ("跨平台对比表",
     "PYNQ-Z2 vs KV260：资源 / 频率 / 吞吐 / 功耗"),
]
for i, (title, detail) in enumerate(plan):
    y = 1.9 + i * 1.42
    add_text(s, title, 10.3, y, 9.4, 0.65,
             font_name="Calibri", font_size=21, bold=True, color=NAVY)
    add_text(s, detail, 10.5, y + 0.65, 9.2, 0.7,
             font_name="Calibri", font_size=19, color=DARK)


# ════════════════════════════════════════════════════════════════════
# Slide 21 — Thank You
# ════════════════════════════════════════════════════════════════════
s = new_slide()
copy_all_freeforms(src_slides[13], s)

add_text(s, "Thank You", 10.0, 4.5, 9.0, 2.5,
         font_name="Cormorant Garamond Bold Italics",
         font_size=108, bold=True, color=NAVY)
add_text(s, "感谢各位老师指导", 10.0, 7.5, 9.0, 1.0,
         font_name="Segoe Print", font_size=32, color=NAVY)
add_text(s, "2026 / 04 / 12", 10.0, 8.8, 9.0, 0.8,
         font_name="Segoe Print", font_size=28, color=NAVY)

# ── 保存 ──────────────────────────────────────────────────────────
prs.save(OUT)
print(f"\n✓ 保存完成：{OUT}")
print(f"  共 {len(prs.slides)} 张幻灯片")